import io
from pptx import Presentation
from supabase import create_client
from core.config import SUPABASE_URL, SUPABASE_SERVICE_KEY

supabase = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)


def _extract_ppt_script(script_text: str) -> str:
    marker = "PPT SCRIPT:"
    if marker in script_text:
        return script_text.split(marker, 1)[1].strip()
    return script_text.strip()


def _add_simple_slides(prs: Presentation, ppt_text: str) -> None:
    """
    Minimal PPT slide creation from the PPT SCRIPT text.
    Safe + simple.
    """
    lines = [ln.strip() for ln in ppt_text.splitlines() if ln.strip()]
    if not lines:
        lines = ["Slide 1: Untitled Lecture"]

    # Use Title and Content layout if available
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    current_title = None
    bullets = []

    def flush_slide():
        nonlocal current_title, bullets
        if not current_title:
            return
        slide = prs.slides.add_slide(layout)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = current_title
        else:
            # fallback if no title placeholder
            tx = slide.shapes.add_textbox(1_000_000, 200_000, 8_000_000, 1_000_000).text_frame
            tx.text = current_title

        # Body
        body = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape is not slide.shapes.title:
                body = shape.text_frame
                break

        if body:
            body.clear()
            for i, b in enumerate(bullets):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = b
                p.level = 0

        current_title = None
        bullets = []

    for ln in lines:
        if ln.lower().startswith("- slide"):
            flush_slide()
            # e.g. "- Slide 2: Title"
            parts = ln.split(":", 1)
            current_title = parts[1].strip() if len(parts) > 1 else ln.replace("- ", "").strip()
        elif ln.startswith("-"):
            bullets.append(ln.lstrip("-").strip())
        else:
            # treat as bullet if we're inside a slide
            if current_title:
                bullets.append(ln)

    flush_slide()


async def generate_pptx_and_upload(lecture_id: str, educator_id: str, script_text: str) -> str:
    ppt_text = _extract_ppt_script(script_text)
    prs = Presentation()
    _add_simple_slides(prs, ppt_text)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    ppt_bytes = buf.read()

    storage_path = f"{educator_id}/{lecture_id}/artifacts/slides.pptx"

    try:
        supabase.storage.from_("lecture-assets").upload(
            storage_path,
            ppt_bytes,
            {
                "content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "upsert": True,
            },
        )
    except Exception as e:
        raise RuntimeError(f"Failed to upload pptx to storage: {e}")

    public = supabase.storage.from_("lecture-assets").get_public_url(storage_path)
    return public
